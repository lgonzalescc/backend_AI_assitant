import os
import re
import json
import logging
import hashlib
import datetime as dt
from typing import List, Tuple, Dict, Optional

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel

import numpy as np
from dataclasses import dataclass
from openai import OpenAI
from fastapi.middleware.cors import CORSMiddleware
import pandas as pd
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import html

app = FastAPI(title="HR Hub Chatbot Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins for testing; restrict in production (e.g., ["http://localhost:3000"])
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@dataclass
class Chunk:
    id: str
    text: str
    source_name: str
    location: str

# Global variables
chunks: List[Chunk] = []
embeddings: np.ndarray = np.array([])

def normalize_whitespace(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"\s+", " ", text)
    return text.strip()

def split_text_into_overlapping_parts(text: str, num_parts: int = 4, overlap_fraction: float = 0.25) -> List[str]:
    words = text.split()
    total_words = len(words)
    if total_words == 0:
        return []
    denominator = num_parts - (num_parts - 1) * overlap_fraction if num_parts > 1 else 1
    part_words = int(total_words / denominator) if denominator > 0 else total_words
    overlap_words = int(part_words * overlap_fraction)
    step = part_words - overlap_words
    parts = []
    start = 0
    for i in range(num_parts):
        end = min(start + part_words, total_words)
        part = ' '.join(words[start:end])
        if part:
            parts.append(part)
        if end >= total_words:
            break
        start += step
    return parts

def read_pdf_chunks(path: str) -> List[Chunk]:
    chunks: List[Chunk] = []
    try:
        reader = PdfReader(path)
        for index, page in enumerate(reader.pages, start=1):
            page_text = normalize_whitespace(page.extract_text() or "")
            if not page_text:
                continue
            parts = split_text_into_overlapping_parts(page_text, 4, 0.25)
            for part_idx, part in enumerate(parts, start=1):
                chunks.append(
                    Chunk(
                        text=part,
                        source_name=os.path.basename(path),
                        location=f"page {index} part {part_idx}",
                        id=f"pdf:{os.path.basename(path)}:p{index}pt{part_idx}",
                    )
                )
        return chunks
    except Exception:
        return []

def read_docx_chunks(path: str) -> List[Chunk]:
    try:
        doc = DocxDocument(path)
        unit_texts: List[str] = []
        unit_locations: List[str] = []
        current_section: Optional[str] = None
        for paragraph_index, paragraph in enumerate(doc.paragraphs, start=1):
            text = normalize_whitespace(paragraph.text)
            if not text:
                continue
            style_name = getattr(paragraph.style, "name", "") or ""
            if style_name.lower().startswith("heading") or style_name.lower().startswith("título"):
                current_section = text
            location = f"section '{current_section}' paragraph {paragraph_index}" if current_section else f"paragraph {paragraph_index}"
            unit_texts.append(text)
            unit_locations.append(location)
        for t_idx, table in enumerate(doc.tables, start=1):
            for r_idx, row in enumerate(table.rows, start=1):
                cells = [normalize_whitespace(cell.text or "") for cell in row.cells]
                cells = [c for c in cells if c]  # remove empty
                if cells:
                    unit_texts.append(" | ".join(cells))
                    unit_locations.append(
                        f"{'section ' + repr(current_section) + ' ' if current_section else ''}table {t_idx} row {r_idx}"
                    )
        return create_overlapped_chunks(unit_texts, unit_locations, os.path.basename(path), "docx", path)
    except Exception:
        pass
    return []

def create_overlapped_chunks(unit_texts: List[str], unit_locations: List[str], basename: str, kind: str, path: str, window_size: int = 4, overlap: int = 1) -> List[Chunk]:
    chunks: List[Chunk] = []
    step = window_size - overlap
    for i in range(0, len(unit_texts), step):
        end = i + window_size
        slice_texts = unit_texts[i:end]
        slice_locations = unit_locations[i:end] if end <= len(unit_locations) else unit_locations[i:]
        if len(slice_texts) == 0:
            continue
        merge_len = len(slice_texts)
        if merge_len < 3 and chunks:
            # merge to last
            last = chunks[-1]
            last.text += ' ' + ' '.join(slice_texts)
            last_loc_end = slice_locations[-1] if slice_locations else ""
            last.location = last.location.rsplit(" to ", 1)[0] + f" to {last_loc_end}"
            # update id
            parts = last.id.split(':')
            if len(parts) == 3 and parts[2].startswith('u'):
                range_part = parts[2][1:]
                start_str, end_str = range_part.split('-')
                start = int(start_str)
                old_end = int(end_str)
                new_end = old_end + merge_len
                last.id = f"{parts[0]}:{parts[1]}:u{start}-{new_end}"
            continue
        chunk_text = ' '.join(slice_texts)
        loc_start = slice_locations[0]
        loc_end = slice_locations[-1]
        chunk_loc = f"{loc_start} to {loc_end}"
        start_idx = i + 1
        end_idx = i + merge_len
        chunk_id = f"{kind}:{basename}:u{start_idx}-{end_idx}"
        chunks.append(
            Chunk(
                text=chunk_text,
                source_name=basename,
                location=chunk_loc,
                id=chunk_id,
            )
        )
    return chunks

def read_pptx_chunks(path: str) -> List[Chunk]:
    chunks: List[Chunk] = []
    try:
        prs = Presentation(path)
        for slide_index, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        txt = "\n".join(p.text for p in shape.text_frame.paragraphs)
                        txt = normalize_whitespace(txt)
                        if txt:
                            texts.append(txt)
                except Exception:
                    continue
            slide_text = normalize_whitespace("\n".join(texts))
            if not slide_text:
                continue
            parts = split_text_into_overlapping_parts(slide_text, 4, 0.25)
            for part_idx, part in enumerate(parts, start=1):
                chunks.append(
                    Chunk(
                        text=part,
                        source_name=os.path.basename(path),
                        location=f"slide {slide_index} part {part_idx}",
                        id=f"pptx:{os.path.basename(path)}:s{slide_index}pt{part_idx}",
                    )
                )
        return chunks
    except Exception:
        return []

def compute_folder_hash(folder_path: str) -> str:
    hash_md5 = hashlib.md5()
    for subdir, _, files in os.walk(folder_path):
        for file in sorted(files):  # sort for consistency
            file_path = os.path.join(subdir, file)
            try:
                with open(file_path, "rb") as f:
                    for chunk in iter(lambda: f.read(4096), b""):
                        hash_md5.update(chunk)
            except Exception:
                pass
    return hash_md5.hexdigest()

def save_chunks_to_excel(chunks: List[Chunk], excel_path: str):
    data = []
    for chunk in chunks:
        data.append({
            'id': chunk.id,
            'text': chunk.text,
            'source_name': chunk.source_name,
            'location': chunk.location
        })
    df = pd.DataFrame(data)
    df.to_excel(excel_path, index=False)

def load_corpus(project_root: str) -> List[Chunk]:
    latest_updates_dir = os.path.join(project_root, 'Latest Updates')
    hash_file = os.path.join(project_root, 'corpus_hash.txt')
    chunks_excel = os.path.join(project_root, 'chunks.xlsx')

    current_hash = compute_folder_hash(latest_updates_dir)
    
    # Check if hash changed
    regenerate = True
    if os.path.exists(hash_file):
        with open(hash_file, 'r') as f:
            saved_hash = f.read().strip()
        if saved_hash == current_hash and os.path.exists(chunks_excel):
            regenerate = False
    
    if not regenerate:
        # Load from excel
        try:
            df = pd.read_excel(chunks_excel)
            chunks = [Chunk(row['id'], row['text'], row['source_name'], row['location']) for _, row in df.iterrows()]
            return chunks
        except Exception:
            regenerate = True
    
    # Regenerate
    supported_ext = {".pdf", ".docx", ".pptx"}
    chunks: List[Chunk] = []
    for dirpath, _, filenames in os.walk(latest_updates_dir):
        for fname in filenames:
            ext = os.path.splitext(fname)[1].lower()
            if ext not in supported_ext:
                continue
            abspath = os.path.join(dirpath, fname)
            try:
                if os.path.getsize(abspath) == 0:
                    continue
            except Exception:
                continue
            new_chunks = []
            try:
                if ext == ".pdf":
                    new_chunks = read_pdf_chunks(abspath)
                elif ext == ".docx":
                    new_chunks = read_docx_chunks(abspath)
                elif ext == ".pptx":
                    new_chunks = read_pptx_chunks(abspath)
            except Exception:
                continue
            chunks.extend(new_chunks)
    
    # Save to excel and update hash
    save_chunks_to_excel(chunks, chunks_excel)
    with open(hash_file, 'w') as f:
        f.write(current_hash)
    
    return chunks

def load_name_to_link(project_root: str) -> Dict[str, str]:
    excel_path = os.path.join(project_root, 'reference', 'SourcesHub_1.xlsx')
    if not os.path.exists(excel_path):
        logging.warning(f"Excel file not found: {excel_path}")
        return {}
    try:
        df = pd.read_excel(excel_path)
    except Exception as e:
        logging.error(f"Error reading Excel: {e}")
        return {}
    name_col = 'NAME DOWNLOAD'
    link_col = 'Link'
    if name_col not in df.columns or link_col not in df.columns:
        logging.warning("Missing columns in Excel")
        return {}
    mapping = {}
    for _, row in df.iterrows():
        name = str(row[name_col]).strip()
        link = str(row[link_col]).strip()
        if name and link:
            mapping[name.lower()] = link
            norm_name = normalize_name(name)  # Asume normalize_name existe; si no, agrégala como en edits previos
            mapping[norm_name] = link
            # Variantes para matching robusto
            variant1 = norm_name.replace('-', '')
            mapping[variant1] = link
            variant2 = name.lower().replace(' ', '').replace('-', '')
            mapping[variant2] = link
            variant3 = re.sub(r'\s*\(\d+\)$|\d+', '', norm_name)
            mapping[variant3] = link
    logging.info(f"Loaded mapping keys from SourcesHub_1.xlsx: {list(mapping.keys())}")
    return mapping

def normalize_name(name: str) -> str:
    name = name.lower()
    name = re.sub(r'\.\w+$', '', name)  # remove extension
    name = re.sub(r'\s*\(\d+\)$', '', name)  # remove (1) etc.
    name = re.sub(r'[^a-z0-9]+', '-', name)  # replace non-alnum with -
    name = name.strip('-')
    return name

def load_sources_reference_map(project_root: str) -> Dict[str, Dict[str, str]]:
    excel_path = os.path.join(project_root, 'reference', 'Sources - P&R Hub 1.xlsx')
    if not os.path.exists(excel_path):
        logging.warning(f"Reference Excel not found: {excel_path}")
        return {}
    try:
        df = pd.read_excel(excel_path, engine='openpyxl')
    except Exception as e:
        logging.error(f"Error reading reference Excel: {e}")
        return {}
    mapping: Dict[str, Dict[str, str]] = {}
    for _, row in df.iterrows():
        name_download = str(row.get('NAME DOWNLOAD', '')).lower().strip()
        if not name_download:
            continue
        title = str(row.get('Document Title', name_download)).strip()
        pr_ref = str(row.get('P&amp;R Hub reference', '')).strip() if pd.notna(row.get('P&amp;R Hub reference')) else ''
        norm_key = normalize_name(name_download)
        mapping[norm_key] = {
            'title': title,
            'ref': pr_ref
        }
        mapping[name_download] = mapping[norm_key]
    logging.info(f"Loaded reference mapping keys: {list(mapping.keys())}")
    return mapping

def evaluate_top_sources(output: str, ranked: List[Tuple[Chunk, float]]) -> List[str]:
    if not os.environ.get("OPENAI_API_KEY", "").strip():
        return []

    context_blocks: List[str] = []
    for c, _ in ranked:
        context_blocks.append(f"ID: {c.id}, Source: {c.source_name}, Content: {c.text}")
    context = "\n\n".join(context_blocks)

    system_msg = "You are an AI evaluator. Given an output answer and chunks with sources, identify the top 3 source documents most crucial to generating the output. Return only the source names in order of relevance, separated by commas. No additional text."

    user_msg = f"Output: {output}\n\nChunks:\n{context}"

    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content": user_msg},
    ]

    try:
        client = OpenAI()
        completion = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0,
            max_tokens=100,
        )
        response = completion.choices[0].message.content.strip()
        top = [s.strip() for s in response.split(',')][:3]
        return top
    except Exception:
        return []

def build_embeddings_index(chunks: List[Chunk]) -> np.ndarray:
    return get_openai_embeddings([c.text for c in chunks])


def rank_chunks(query: str, embeddings: np.ndarray, chunks: List[Chunk], top_k: int = 25) -> List[Tuple[Chunk, float]]:
    if not query.strip() or embeddings.shape[0] == 0:
        return []
    q_vec = get_openai_embeddings([query])[0]
    sims = cosine_sim(q_vec, embeddings)
    scored = [(idx, float(s)) for idx, s in enumerate(sims) if s > 0.0]
    scored.sort(key=lambda x: x[1], reverse=True)
    out = []
    for idx, s in scored[:top_k]:
        out.append((chunks[idx], s))
    return out

def cosine_sim(query_vec: np.ndarray, matrix: np.ndarray) -> np.ndarray:
    query_vec = query_vec / np.linalg.norm(query_vec)
    matrix = matrix / np.linalg.norm(matrix, axis=1, keepdims=True)
    return np.dot(matrix, query_vec)


def get_openai_embeddings(texts: List[str]) -> np.ndarray:
    client = OpenAI()
    resp = client.embeddings.create(
        model="text-embedding-3-small",
        input=texts
    )
    vectors = [d.embedding for d in resp.data]
    return np.array(vectors, dtype=np.float32)


# Add other functions from the snippet, like extract_relevant_sentences, format_sources_lines, call_openai_generate, etc.

def call_openai_generate(query: str, ranked: List[Tuple[Chunk, float]], max_sentences: int = 5, custom_system_msg: Optional[str] = None, name_to_link: Dict[str, str] = {}) -> Optional[str]:
    selected = ranked
    context_blocks: List[str] = []
    for c, _ in selected:
        link = name_to_link.get(c.source_name, '')
        context_blocks.append(
            f"[ID: {c.id}]\nLink: {link if link else 'N/A'}\nFile: {c.source_name}\nLocation: {c.location}\nContent: {c.text}"
        )
    context = "\n\n---\n\n" + "\n\n---\n\n".join(context_blocks) if context_blocks else ""

    base_msg = custom_system_msg if custom_system_msg is not None else "You are a helpful assistant that answer user questions in a detailed and comprehensive way."

    system_msg = base_msg + "\n\nAlways base your answer strictly on the provided context.\nProvide a detailed and comprehensive answer to the question."

    user_msg = (
        f"Question: {query}\n\nContext:{context}"
    )

    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content": user_msg},
    ]
    full_input = json.dumps(messages, ensure_ascii=False)

    try:
        api_key = os.environ.get("OPENAI_API_KEY", "").strip()
        if not api_key:
            return None
        client = OpenAI(api_key=api_key)
        try:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=messages,
                temperature=0,
                max_tokens=5000,
            )
            return (resp.choices[0].message.content or "").strip()
        except Exception:
            try:
                resp2 = client.responses.create(
                    model="gpt-4o-mini",
                    input=messages,
                    temperature=0,
                    max_output_tokens=5000,
                )
                if hasattr(resp2, "output") and resp2.output and hasattr(resp2.output[0], "content"):
                    parts = resp2.output[0].content
                    if parts and hasattr(parts[0], "text"):
                        return (parts[0].text or "").strip()
            except Exception:
                return None
    except Exception:
        return None

    return None

@app.on_event("startup")
def startup():
    global chunks, embeddings
    project_root = os.path.dirname(os.path.abspath(__file__))
    chunks = load_corpus(project_root)
    embeddings = build_embeddings_index(chunks)

class QueryRequest(BaseModel):
    query: str
    top_k: int = 200
    max_sentences: int = 5

@app.post("/query")
def process_query(request: QueryRequest):
    try:
        project_root = os.path.dirname(os.path.abspath(__file__))
        name_to_link = load_name_to_link(project_root)
        ranked = rank_chunks(request.query, embeddings, chunks, request.top_k)
        answer = call_openai_generate(request.query, ranked, request.max_sentences, name_to_link=name_to_link)
        if not answer:
            answer = "Fallback answer"
        
        top_sources = evaluate_top_sources(answer, ranked)
        
        sources_ref_map = load_sources_reference_map(project_root)
        sources_table = []
        for source_name in top_sources:
            norm_source = normalize_name(source_name)
            link = name_to_link.get(norm_source, name_to_link.get(source_name.lower(), 'N/A'))
            if link == 'N/A':
                # Enhanced fuzzy match
                variant = norm_source.replace('-', '').replace(' ', '')
                link = name_to_link.get(variant, 'N/A')
            if link == 'N/A':
                variant = re.sub(r'\s*\(\d+\)$|\d+', '', norm_source)
                link = name_to_link.get(variant, 'N/A')
            if link == 'N/A':
                # Partial containment fallback
                for key in name_to_link:
                    if variant in key or key in variant:
                        link = name_to_link[key]
                        break
            # Get title from ref map as before
            meta = sources_ref_map.get(norm_source, sources_ref_map.get(source_name.lower(), {}))
            title = meta.get('title', source_name)
            sources_table.append({
                'title': title,
                'link': link
            })
        
        return {
            "answer": answer,
            "sources_table": sources_table
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
